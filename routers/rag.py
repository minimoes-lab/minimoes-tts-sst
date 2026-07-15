# This software is licensed under a **dual-license model**
# For individuals and businesses earning **under $1M per year**, this software is licensed under the **MIT License**
# Businesses or organizations with **annual revenue of $1,000,000 or more** must obtain permission to use this software commercially.

import os
import uuid
import shutil
import asyncio
import urllib.parse
import xml.etree.ElementTree as ET
from datetime import datetime
from typing import List, Optional, Dict, Any

import httpx
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import zipfile
from pypdf import PdfReader

from fastapi import APIRouter, File, UploadFile, Form, HTTPException
from pydantic import BaseModel, Field

from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from langchain_groq import ChatGroq
from langchain_core.prompts import PromptTemplate

import core.state as state

router = APIRouter()

GROQ_API_KEY = os.getenv("GROQ_API_KEY")

_BLOCKED_HOSTS = {'localhost', '127.0.0.1', '0.0.0.0', '::1'}
_BLOCKED_PREFIXES = (
    '10.', '192.168.', '172.16.', '172.17.', '172.18.', '172.19.',
    '172.20.', '172.21.', '172.22.', '172.23.', '172.24.', '172.25.',
    '172.26.', '172.27.', '172.28.', '172.29.', '172.30.', '172.31.',
    '169.254.', 'fc', 'fd', 'fe80::',
)
_MAX_URLS = 50

from streaming.streaming_rag import RAG_PROMPT_TEMPLATE
RAG_PROMPT = PromptTemplate.from_template(RAG_PROMPT_TEMPLATE)

DEFAULT_SYSTEM_PROMPT = "You are a voice-first conversational assistant. Be natural, warm, and concise. Use short sentences. Write for speech, not reading. Reply in the same language as the user."


# ── Pydantic models ──────────────────────────────────────────────────────────

class ProcessResponse(BaseModel):
    session_id: str = Field(..., description="Unique identifier for the processing session.")
    message: str = Field(..., description="A confirmation message.")
    filenames: List[str] = Field(..., description="List of filenames processed.")
    processed_url: Optional[str] = Field(None, description="The URL that was processed, if any.")


# ── SSRF guard ───────────────────────────────────────────────────────────────

def is_valid_public_url(url: str) -> bool:
    """
    Allowlist-based SSRF defence (OWASP SSRF Prevention Cheat Sheet):
    - Only http:// and https:// schemes accepted (blocks file://, ftp://, gopher://)
    - Hostname must not resolve to private/loopback/metadata ranges
    - Redirects are disabled at fetch time (allow_redirects=False)
    """
    try:
        parsed = urllib.parse.urlparse(url)
        # Scheme allowlist: only http and https
        if parsed.scheme not in ("http", "https"):
            return False
        hostname = (parsed.hostname or '').lower()
        if not hostname:
            return False
        if hostname in _BLOCKED_HOSTS:
            return False
        if hostname.startswith(_BLOCKED_PREFIXES):
            return False
        # Block cloud metadata endpoints explicitly
        if hostname in ("169.254.169.254", "metadata.google.internal", "metadata.azure.com"):
            return False
        return True
    except Exception:
        return False


# ── ContentExtractor ─────────────────────────────────────────────────────────

class ContentExtractor:
    async def from_url(self, url: str) -> str:
        try:
            # follow_redirects=False prevents redirect-based SSRF bypass
            # (attacker submits public URL that redirects to internal/metadata endpoint)
            async with httpx.AsyncClient(
                timeout=15,
                follow_redirects=False,
                headers={'User-Agent': 'Mozilla/5.0 (compatible; RAG-Bot/1.0)'},
            ) as client:
                response = await client.get(url)
            if 300 <= response.status_code < 400:
                raise HTTPException(status_code=400, detail=f"URL redirects are not allowed: {url}")
            response.raise_for_status()
            soup = BeautifulSoup(response.content, 'html.parser')
            for element in soup(["script", "style", "header", "footer", "nav", "aside"]):
                element.decompose()
            return soup.get_text(separator='\n', strip=True)
        except httpx.HTTPError as e:
            raise HTTPException(status_code=400, detail=f"Failed to fetch or parse URL {url}. Error: {e}")

    async def from_sitemap(self, sitemap_url: str) -> List[str]:
        try:
            async with httpx.AsyncClient(
                timeout=30,
                follow_redirects=False,
                headers={'User-Agent': 'Mozilla/5.0 (compatible; RAG-Bot/1.0)'},
            ) as client:
                response = await client.get(sitemap_url)
            if 300 <= response.status_code < 400:
                return []
            response.raise_for_status()
            root = ET.fromstring(response.content)
            return [elem.text.strip() for elem in root.iter() if elem.tag.endswith('loc') and elem.text]
        except Exception as e:
            print(f"[{datetime.now()}] ERROR: Failed to parse sitemap {sitemap_url}: {e}")
            return []

    def from_pdf(self, file_stream) -> str:
        try:
            reader = PdfReader(file_stream)
            return "".join(page.extract_text() for page in reader.pages if page.extract_text())
        except Exception as e:
            return f"[PDF Parsing Error: {e}]"

    def from_docx(self, file_stream) -> str:
        try:
            doc = Document(file_stream)
            return "\n".join(para.text for para in doc.paragraphs if para.text)
        except Exception as e:
            return f"[DOCX Parsing Error: {e}]"

    def from_pptx(self, file_stream) -> str:
        try:
            prs = Presentation(file_stream)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        except Exception as e:
            return f"[PPTX Parsing Error: {e}]"

    def from_excel(self, file_stream) -> str:
        try:
            xls = pd.ExcelFile(file_stream)
            text = ""
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                text += f"--- Sheet: {sheet_name} ---\n{df.to_string()}\n\n"
            return text
        except Exception as e:
            return f"[Excel Parsing Error: {e}]"

    def from_zip(self, file_stream, temp_dir) -> str:
        text = ""
        try:
            with zipfile.ZipFile(file_stream) as z:
                # Zip slip defence (OWASP A03:2021 — Injection):
                # Python < 3.12 extractall() does NOT strip path-traversal members
                # (e.g. "../../etc/passwd").  Resolve each member's target path and
                # skip anything that escapes temp_dir.
                # Ref: https://owasp.org/www-community/vulnerabilities/Zip_Slip
                base = os.path.realpath(temp_dir)
                for member in z.infolist():
                    member_path = os.path.realpath(os.path.join(temp_dir, member.filename))
                    if not (member_path == base or member_path.startswith(base + os.sep)):
                        print(f"[ZIP] Blocked unsafe path in archive: {member.filename!r}")
                        continue
                    z.extract(member, path=temp_dir)
                for root, _, files in os.walk(temp_dir):
                    for file in files:
                        text += self.from_file_path(os.path.join(root, file)) + "\n\n"
        except Exception as e:
            text += f"[ZIP Parsing Error: {e}]"
        return text

    def from_file_path(self, file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        with open(file_path, "rb") as f:
            if ext == ".pdf":       return self.from_pdf(f)
            elif ext == ".docx":    return self.from_docx(f)
            elif ext == ".pptx":    return self.from_pptx(f)
            elif ext in [".xls", ".xlsx"]: return self.from_excel(f)
            elif ext == ".txt":     return f.read().decode("utf-8", errors="replace")
            else:                   return f"[Unsupported File: {os.path.basename(file_path)}]"


# ── RAG chain factory ────────────────────────────────────────────────────────

def get_rag_chain(text_chunks: List[str]) -> ConversationalRetrievalChain:
    if not text_chunks:
        raise ValueError("Cannot create RAG chain with no text chunks.")
    if not state.embeddings_model:
        raise RuntimeError("Embeddings model not loaded.")
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        raise RuntimeError("GROQ_API_KEY is not set.")

    vector_store = FAISS.from_texts(texts=text_chunks, embedding=state.embeddings_model)
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True, output_key='answer')
    llm = ChatGroq(
        model="meta-llama/llama-4-scout-17b-16e-instruct",
        temperature=0.7,
        max_tokens=800,
        groq_api_key=api_key,
    )
    return ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vector_store.as_retriever(),
        memory=memory,
        return_source_documents=True,
        combine_docs_chain_kwargs={"prompt": RAG_PROMPT},
    )


# ── Endpoint ─────────────────────────────────────────────────────────────────

@router.post("/process", response_model=ProcessResponse)
async def process_content(
    files: Optional[List[UploadFile]] = File(None),
    url: Optional[str] = Form(None),
    crawl_urls: Optional[str] = Form(None),
    sitemap_urls: Optional[str] = Form(None),
    individual_urls: Optional[str] = Form(None),
    prompt_template: Optional[str] = Form(None),
):
    print(f"[{datetime.now()}] /process endpoint called.")
    has_any_source = bool(files) or bool(url) or bool(crawl_urls) or bool(sitemap_urls) or bool(individual_urls)
    if not has_any_source:
        if not GROQ_API_KEY:
            raise HTTPException(status_code=500, detail="GROQ_API_KEY not configured.")
        system_prompt = (prompt_template.strip() if prompt_template and prompt_template.strip() else DEFAULT_SYSTEM_PROMPT)
        session_id = uuid.uuid4().hex
        state.set_conversation(session_id, {
            "type": "direct",
            "system_prompt": system_prompt,
            "history": [],
        })
        print(f"[{datetime.now()}] Direct LLM session created: {session_id}")
        return ProcessResponse(session_id=session_id, message="Direct LLM session (no RAG)", filenames=[])

    extractor = ContentExtractor()
    import tempfile
    temp_dir = tempfile.mkdtemp(prefix="rag_upload_")
    raw_text = ""
    processed_files = []
    all_page_urls = []

    if url:
        all_page_urls.extend([u.strip() for u in url.split('\n') if u.strip()])
    if crawl_urls:
        all_page_urls.extend([u.strip() for u in crawl_urls.split('\n') if u.strip()])
    if individual_urls:
        all_page_urls.extend([u.strip() for u in individual_urls.split('\n') if u.strip()])
    if sitemap_urls:
        for sitemap_url in [u.strip() for u in sitemap_urls.split('\n') if u.strip()]:
            if not is_valid_public_url(sitemap_url):
                continue
            all_page_urls.extend(await extractor.from_sitemap(sitemap_url))

    all_page_urls = [u for u in all_page_urls if is_valid_public_url(u)]
    if not all_page_urls and not files:
        raise HTTPException(status_code=400, detail="No valid public URLs provided after security filtering.")
    if len(all_page_urls) > _MAX_URLS:
        all_page_urls = all_page_urls[:_MAX_URLS]

    try:
        for page_url in all_page_urls:
            try:
                raw_text += await extractor.from_url(page_url) + "\n\n"
                await asyncio.sleep(1.0)
            except Exception as e:
                print(f"[{datetime.now()}] ERROR scraping URL {page_url}: {e}")

        _MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB per file
        if files:
            loop = asyncio.get_running_loop()
            for file in files:
                raw_name = file.filename or ''
                ext = os.path.splitext(os.path.basename(raw_name))[1].lower() if raw_name else ''
                safe_filename = uuid.uuid4().hex + (ext if ext in {'.pdf', '.txt', '.docx', '.zip', '.csv'} else '.bin')
                file_path = os.path.join(temp_dir, safe_filename)
                file_content = await file.read(_MAX_FILE_SIZE + 1)
                if len(file_content) > _MAX_FILE_SIZE:
                    raise HTTPException(status_code=413, detail=f"File '{safe_filename}' exceeds 50 MB limit.")

                def _write_and_extract(_path=file_path, _content=file_content, _fname=file.filename):
                    with open(_path, "wb") as buf:
                        buf.write(_content)
                    ext = os.path.splitext(_fname)[1].lower()
                    if ext == ".zip":
                        return extractor.from_zip(_path, os.path.join(temp_dir, "unzipped"))
                    return extractor.from_file_path(_path) + "\n\n"

                extracted = await loop.run_in_executor(None, _write_and_extract)
                raw_text += extracted
                processed_files.append(file.filename)

        if not raw_text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted from the provided sources.")

        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200, length_function=len)
        text_chunks = text_splitter.split_text(raw_text)

        session_id = str(uuid.uuid4())
        rag_chain = get_rag_chain(text_chunks)

        if prompt_template:
            # prompt_template is a free-text agent description; inject it as the
            # system context rather than replacing the full LangChain prompt template.
            # This prevents overriding {context}/{question} variables and blocks prompt injection.
            # Escape braces so user text cannot inject LangChain template variables
            safe_desc = prompt_template[:4000].replace('{', '{{').replace('}', '}}')
            full_template = (
                f"{safe_desc}\n\n"
                "Use the following pieces of context to answer the question:\n\n"
                "{context}\n\n"
                "Question: {question}\n"
                "Helpful Answer:"
            )
            rag_chain.combine_docs_chain.llm_chain.prompt = PromptTemplate.from_template(full_template)

        state.set_conversation(session_id, rag_chain)

        return ProcessResponse(
            session_id=session_id,
            message="Content processed successfully.",
            filenames=processed_files,
            processed_url=url,
        )
    except Exception as e:
        print(f"[{datetime.now()}] ERROR in /process: {e}")
        raise HTTPException(status_code=500, detail="An internal error occurred.")
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

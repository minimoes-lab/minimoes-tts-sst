

# This software is licensed under a **dual-license model**
# For individuals and businesses earning **under $1M per year**, this software is licensed under the **MIT License**
# Businesses or organizations with **annual revenue of $1,000,000 or more** must obtain permission to use this software commercially.
from fastapi import Request  # Capital "R" for FastAPI
from fastapi.responses import JSONResponse
from datetime import datetime
import time
import io
import base64
import os
os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")


import uuid
import shutil
import uvicorn
import requests
import re
import tempfile
import threading
from bs4 import BeautifulSoup
from fastapi import FastAPI, File, UploadFile, Form, HTTPException, BackgroundTasks, WebSocket, WebSocketDisconnect
from fastapi.responses import JSONResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field
from typing import List, Optional, Dict, Any

# --- Core Machine Learning & NLP Imports ---
import torch
import numpy as np
import pandas as pd
from docx import Document
from pptx import Presentation
import zipfile
from PyPDF2 import PdfReader
import scipy.io.wavfile as wavfile
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from langchain_groq import ChatGroq
from langchain_core.prompts import PromptTemplate

# =====================================================================================
# 1. API Initialization & Configuration
# =====================================================================================


from utils.generate_face_shapes import generate_facial_data_from_bytes
from utils.model.model import load_model
from utils.config import config, get_blendshape_names, blendshapes_to_named_frames
import multiprocessing

from streaming.qwen_tts_worker import QwenTTSWorker
app = FastAPI(
    title="Intelligent Document & Web API",
    description="A high-quality API for querying documents and websites using a RAG pipeline with Groq, and generating speech with Qwen3-TTS.",
    version="2.0.1"
)


if __name__ == '__main__' or __name__.startswith("api"):
    try:
        multiprocessing.set_start_method('spawn', force=True)
    except RuntimeError:
        pass

device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
print("Activated device:", device)
BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# Join that with the filename
model_path = "utils/model/model.pth"

@app.get("/")
def root():
    return {"status": "ok"}

@app.get("/health")
def health():
    return {"status": "healthy"}


_tts_worker_for_speakers: Optional[QwenTTSWorker] = None

_tts_reference_audio_path: Optional[str] = os.getenv("TTS_REF_AUDIO_PATH")
_tts_reference_text: Optional[str] = os.getenv("TTS_REF_TEXT")


@app.post("/tts/reference_audio")
async def set_tts_reference_audio(
    audio: UploadFile = File(...),
    text: str = Form(...),
):
    global _tts_reference_audio_path, _tts_reference_text, _tts_worker_for_speakers

    if audio is None:
        raise HTTPException(status_code=400, detail="Missing audio file")
    if not text or not text.strip():
        raise HTTPException(status_code=400, detail="Missing reference text")

    filename = (audio.filename or "").lower()
    if filename and not filename.endswith(".wav"):
        raise HTTPException(status_code=400, detail="Reference audio must be a .wav file")

    ref_dir = os.path.join(BASE_DIR, "tts_reference")
    os.makedirs(ref_dir, exist_ok=True)
    ref_path = os.path.join(ref_dir, f"ref_{uuid.uuid4().hex}.wav")

    try:
        content = await audio.read()
        if not content:
            raise HTTPException(status_code=400, detail="Empty audio file")
        with open(ref_path, "wb") as f:
            f.write(content)
    finally:
        try:
            await audio.close()
        except Exception:
            pass

    _tts_reference_audio_path = ref_path
    _tts_reference_text = text.strip()

    # Invalidate cached worker so /tts/speakers reflects new reference immediately
    _tts_worker_for_speakers = None

    return {
        "status": "ok",
        "reference_configured": True,
        "reference_audio_path": _tts_reference_audio_path,
    }


@app.get("/tts/speakers")
def get_tts_speakers():
    global _tts_worker_for_speakers
    global _tts_reference_audio_path, _tts_reference_text
    tts_ref_audio = _tts_reference_audio_path
    tts_ref_text = _tts_reference_text

    if _tts_worker_for_speakers is None:
        device_str = "cuda" if torch.cuda.is_available() else "cpu"
        try:
            _tts_worker_for_speakers = QwenTTSWorker(
                device=device_str,
                use_qwen3=True,
                reference_audio_path=tts_ref_audio,
                reference_text=tts_ref_text,
            )
        except Exception as e:
            return {
                "speakers": [],
                "default_speaker": None,
                "tts_mode": "base_voice_clone",
                "reference_configured": bool(tts_ref_audio) and bool(tts_ref_text),
                "error": str(e),
            }

    speakers = getattr(_tts_worker_for_speakers, "speakers", None) or []
    default_speaker = getattr(_tts_worker_for_speakers, "default_speaker", None)
    return {
        "speakers": speakers,
        "default_speaker": default_speaker,
        "tts_mode": "base_voice_clone",
        "reference_configured": bool(tts_ref_audio) and bool(tts_ref_text)
    }

# Load blendshape model at startup (moved to startup event)
@app.post("/audio_to_blendshapes")
async def audio_to_blendshapes_route(request: Request):

    audio_bytes = await request.body()
    generated_facial_data = generate_facial_data_from_bytes(audio_bytes, blendshape_model, device, config)
    generated_facial_data_list = generated_facial_data.tolist() if isinstance(generated_facial_data, np.ndarray) else generated_facial_data
    
    frames = blendshapes_to_named_frames(generated_facial_data_list)
    
    return JSONResponse(content={
        'frame_rate': config['frame_rate'],
        'total_frames': len(frames),
        'frames': frames,
        'mapping': get_blendshape_names(),
    })


STATIC_AUDIO_DIR = "/app/generated_audio"



# This software is licensed under a **dual-license model**
# For individuals and businesses earning **under $1M per year**, this software is licensed under the **MIT License**









# --- Configuration & Global State ---
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "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")
if not GROQ_API_KEY:
    raise ValueError("GROQ_API_KEY is not set. Please set it as an environment variable.")

conversations: Dict[str, ConversationalRetrievalChain] = {}
embeddings_model = None

# =====================================================================================
# 2. High-Quality RAG Prompt Template
# =====================================================================================

PROFESSIONAL_RAG_PROMPT_TEMPLATE = """
You are a voice-first conversational assistant. Your goal is to sound natural, warm, and human while staying strictly grounded in the provided context.

**Hard rules (must follow):**
1. **Strict grounding:** Use ONLY the provided context and chat history. If the answer is not in the context, you MUST say exactly: "I am sorry, but the information required to answer your question is not available in the provided documents." Do not guess.
2. **Spoken style (TTS-friendly):** Write for speech, not for reading.
   - Use short sentences.
   - Avoid long paragraphs, markdown, and long lists.
   - Prefer simple punctuation to create rhythm (commas, periods, occasional "...").
3. **Brevity:** Default to 1-3 short sentences total.
4. **Emotional attunement:** If the user expresses emotion (stress, frustration, sadness, excitement), acknowledge it briefly in a calm, supportive way.
5. **Keep the conversation moving:** End with ONE short follow-up question to clarify or advance the dialogue.
6. **Language:** Reply in the same language as the user's question.

**Context:**
{context}

**Chat History:**
{chat_history}

**Question:**
{question}

**Answer:**
"""
RAG_PROMPT = PromptTemplate.from_template(PROFESSIONAL_RAG_PROMPT_TEMPLATE)

# =====================================================================================
# 3. Model Loading (at Startup)
# =====================================================================================
from fastapi.responses import PlainTextResponse

@app.on_event("startup")
async def load_models():
    """Load heavy ML models once when the API starts up."""
    global embeddings_model, blendshape_model

    print(f"[{datetime.now()}] Starting model loading...")
    
    print(f"[{datetime.now()}] Loading HuggingFace embeddings model...")
    embeddings_start_time = time.time()
    embeddings_model = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2", model_kwargs={'device': 'cuda' if torch.cuda.is_available() else 'cpu'})
    embeddings_end_time = time.time()
    print(f"[{datetime.now()}] Embeddings model loaded successfully in {embeddings_end_time - embeddings_start_time:.2f} seconds.")
    
    print(f"--- ATTEMPTING TO LOAD: {model_path} ---")
    blendshape_model = load_model(model_path, config, device)
    print(f"DEBUG: Absolute path is: {model_path}")
    
    print(f"[{datetime.now()}] Model loading complete. Using Qwen3-TTS for speech generation.")

# =====================================================================================
# 4. Pydantic Models for API Data Validation
# =====================================================================================

class ProcessResponse(BaseModel):
    session_id: str = Field(..., description="Unique identifier for the processing session.")
    message: str = Field(..., description="A confirmation message.")
    filenames: List[str] = Field(..., description="List of filenames processed.")
    processed_url: Optional[str] = Field(None, description="The URL that was processed, if any.")

class QueryRequest(BaseModel):
    session_id: str = Field(..., description="The session ID obtained from the /process endpoint.")
    question: str = Field(..., description="The question to ask the RAG system.")
    voice_preset: Optional[str] = Field(None, description="Optional voice preset for TTS.")

class SourceDocument(BaseModel):
    page_content: str = Field(..., description="The text content of the source chunk.")
    metadata: Dict[str, Any] = Field(..., description="Metadata about the source, like file name or page number.")

class QueryResponse(BaseModel):
    answer: str = Field(..., description="The generated answer from the RAG model.")
    source_documents: List[SourceDocument] = Field(..., description="List of source document chunks used for the answer.")
    audio_base64: Optional[str] = Field(None, description="Base64 encoded audio content.")

# =====================================================================================
# 5. Helper Functions for Document and Web Parsing
# =====================================================================================

class ContentExtractor:
    """A centralized class for extracting text from various sources."""
    def from_url(self, url: str) -> str:
        try:
            print(f"[{datetime.now()}] Extracting content from URL: {url}")
            response = requests.get(url, timeout=15, headers={'User-Agent': 'Mozilla/5.0'})
            response.raise_for_status()
            soup = BeautifulSoup(response.content, 'html.parser')
            for element in soup(["script", "style", "header", "footer", "nav", "aside"]):
                element.decompose()
            text = soup.get_text(separator='\n', strip=True)
            print(f"[{datetime.now()}] URL content extracted successfully.")
            return text
        except requests.RequestException as e:
            print(f"[{datetime.now()}] ERROR: Failed to fetch or parse URL {url}. Error: {e}")
            raise HTTPException(status_code=400, detail=f"Failed to fetch or parse URL {url}. Error: {e}")

    def from_pdf(self, file_stream) -> str:
        try:
            print(f"[{datetime.now()}] Extracting content from PDF.")
            reader = PdfReader(file_stream)
            text = "".join(page.extract_text() for page in reader.pages if page.extract_text())
            print(f"[{datetime.now()}] PDF content extracted successfully.")
            return text
        except Exception as e:
            print(f"[{datetime.now()}] ERROR: PDF Parsing Error: {e}")
            return f"[PDF Parsing Error: {e}]"

    def from_docx(self, file_stream) -> str:
        try:
            print(f"[{datetime.now()}] Extracting content from DOCX.")
            doc = docx.Document(file_stream)
            text = "\n".join(para.text for para in doc.paragraphs if para.text)
            print(f"[{datetime.now()}] DOCX content extracted successfully.")
            return text
        except Exception as e:
            print(f"[{datetime.now()}] ERROR: DOCX Parsing Error: {e}")
            return f"[DOCX Parsing Error: {e}]"

    def from_pptx(self, file_stream) -> str:
        try:
            print(f"[{datetime.now()}] Extracting content from PPTX.")
            prs = Presentation(file_stream)
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
            print(f"[{datetime.now()}] PPTX content extracted successfully.")
            return text
        except Exception as e:
            print(f"[{datetime.now()}] ERROR: PPTX Parsing Error: {e}")
            return f"[PPTX Parsing Error: {e}]"

    def from_excel(self, file_stream) -> str:
        try:
            print(f"[{datetime.now()}] Extracting content from Excel.")
            xls = pd.ExcelFile(file_stream)
            text = ""
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                text += f"--- Sheet: {sheet_name} ---\n{df.to_string()}\n\n"
            print(f"[{datetime.now()}] Excel content extracted successfully.")
            return text
        except Exception as e:
            print(f"[{datetime.now()}] ERROR: Excel Parsing Error: {e}")
            return f"[Excel Parsing Error: {e}]"

    def from_zip(self, file_stream, temp_dir) -> str:
        print(f"[{datetime.now()}] Extracting content from ZIP.")
        text = ""
        try:
            with zipfile.ZipFile(file_stream) as z:
                z.extractall(path=temp_dir)
                for root, _, files in os.walk(temp_dir):
                    for file in files:
                        file_path = os.path.join(root, file)
                        text += self.from_file_path(file_path) + "\n\n"
            print(f"[{datetime.now()}] ZIP content extracted successfully.")
        except Exception as e:
            print(f"[{datetime.now()}] ERROR: ZIP Parsing Error: {e}")
            text += f"[ZIP Parsing Error: {e}]"
        return text

    def from_file_path(self, file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        print(f"[{datetime.now()}] Processing file: {os.path.basename(file_path)} with extension {ext}")
        with open(file_path, "rb") as f:
            if ext == ".pdf": return self.from_pdf(f)
            elif ext == ".docx": return self.from_docx(f)
            elif ext == ".pptx": return self.from_pptx(f)
            elif ext in [".xls", ".xlsx"]: return self.from_excel(f)
            else:
                print(f"[{datetime.now()}] INFO: Unsupported File in ZIP: {os.path.basename(file_path)}")
                return f"[Unsupported File in ZIP: {os.path.basename(file_path)}]"

# =====================================================================================
# 6. Core RAG and TTS Logic
# =====================================================================================
def get_rag_chain(text_chunks: List[str]) -> ConversationalRetrievalChain:
    print(f"[{datetime.now()}] Creating RAG chain...")
    if not text_chunks:
        raise ValueError("Cannot create RAG chain with no text chunks.")
    if not embeddings_model:
        raise RuntimeError("Embeddings model not loaded.")
    try:
        vector_store_start_time = time.time()
        vector_store = FAISS.from_texts(texts=text_chunks, embedding=embeddings_model)
        vector_store_end_time = time.time()
        print(f"[{datetime.now()}] FAISS vector store created in {vector_store_end_time - vector_store_start_time:.2f} seconds.")

        memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True, output_key='answer')
        llm = ChatGroq(
                     model="meta-llama/llama-4-scout-17b-16e-instruct",
                     temperature=0.7,
                     max_tokens=800,
                     groq_api_key=GROQ_API_KEY
        )

        rag_chain = ConversationalRetrievalChain.from_llm(
            llm=llm,
            retriever=vector_store.as_retriever(),
            memory=memory,
            return_source_documents=True,
            combine_docs_chain_kwargs={"prompt": RAG_PROMPT}
        )
        print(f"[{datetime.now()}] RAG chain created successfully.")
        return rag_chain
    except Exception as e:
        print(f"[{datetime.now()}] ERROR: Failed to create RAG chain: {e}")
        raise RuntimeError(f"Failed to create RAG chain: {e}")

# Add this import near the top of the file
from collections.abc import Mapping

def generate_speech_content_base64(text: str, voice_preset: Optional[str]) -> Optional[str]:
    """Generate speech using Qwen3-TTS and return base64 encoded audio."""
    print(f"[{datetime.now()}] [Qwen3-TTS] Starting speech generation (Base64 encoding).")
    if not text:
        print(f"[{datetime.now()}] [Qwen3-TTS] No text provided. Skipping.")
        return None

    try:
        # Use Qwen TTS worker
        from streaming.qwen_tts_worker import QwenTTSWorker
        
        # Use CPU/CUDA based on availability
        device = "cuda" if torch.cuda.is_available() else "cpu"
        worker = QwenTTSWorker(device=device, use_qwen3=True)
        
        # Generate audio synchronously
        result = worker._generate_audio_sync(text, voice_preset)
        
        if result is None:
            print(f"[{datetime.now()}] [Qwen3-TTS] Failed to generate audio.")
            return None
        
        audio_np, wav_bytes = result
        
        # Encode to base64
        audio_base64 = base64.b64encode(wav_bytes).decode("utf-8")
        
        print(f"[{datetime.now()}] [Qwen3-TTS] Audio generated successfully. Size: {len(wav_bytes)} bytes")
        return audio_base64
        
    except Exception as e:
        print(f"[{datetime.now()}] [Qwen3-TTS] ERROR: {e}")
        import traceback
        traceback.print_exc()
        return None


# =====================================================================================
# 7. API Endpoints
# =====================================================================================

# The /audio mount is technically no longer needed for TTS audio, but kept for static files
# if other non-TTS static files are needed.
STATIC_AUDIO_DIR = "generated_audio"
os.makedirs(STATIC_AUDIO_DIR, exist_ok=True)
app.mount("/audio", StaticFiles(directory=STATIC_AUDIO_DIR), name="audio")

@app.post("/process", response_model=ProcessResponse)
async def process_content(
    files: Optional[List[UploadFile]] = File(None, description="A list of documents to process."),
    url: Optional[str] = Form(None, description="A URL to a website to scrape for text."),
    prompt_template: Optional[str] = Form(None, description="Optional custom RAG prompt template. Must include {context}, {chat_history}, and {question}.")
):
    print(f"[{datetime.now()}] /process endpoint called.")
    if not files and not url:
        raise HTTPException(status_code=400, detail="Please provide at least one document or a URL.")
    
    extractor = ContentExtractor()
    temp_dir = f"temp_{uuid.uuid4().hex}"
    os.makedirs(temp_dir)
    raw_text = ""
    processed_files = []
    
    try:
        process_start_time = time.time()
        if url:
            print(f"[{datetime.now()}] Processing URL: {url}")
            raw_text += extractor.from_url(url) + "\n\n"
        
        if files:
            print(f"[{datetime.now()}] Processing files: {[f.filename for f in files]}")
            for file in files:
                file_path = os.path.join(temp_dir, file.filename)
                with open(file_path, "wb") as buffer:
                    shutil.copyfileobj(file.file, buffer)
                ext = os.path.splitext(file.filename)[1].lower()
                if ext == ".zip":
                    raw_text += extractor.from_zip(file_path, os.path.join(temp_dir, "unzipped"))
                else:
                    raw_text += extractor.from_file_path(file_path) + "\n\n"
                processed_files.append(file.filename)
        
        if not raw_text.strip():
            print(f"[{datetime.now()}] ERROR: No text extracted from sources.")
            raise HTTPException(status_code=400, detail="No text could be extracted from the provided sources.")
        
        print(f"[{datetime.now()}] Splitting text into chunks...")
        text_splitter_start_time = time.time()
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200, length_function=len)
        text_chunks = text_splitter.split_text(raw_text)
        text_splitter_end_time = time.time()
        print(f"[{datetime.now()}] Text split into {len(text_chunks)} chunks in {text_splitter_end_time - text_splitter_start_time:.2f} seconds.")
        
        session_id = str(uuid.uuid4())
        print(f"[{datetime.now()}] Creating RAG chain for session: {session_id}")
        conversations[session_id] = get_rag_chain(text_chunks)
        
        if prompt_template:
            print(f"[{datetime.now()}] Applying custom prompt template.")
            conversations[session_id].combine_docs_chain.llm_chain.prompt = PromptTemplate.from_template(prompt_template)
        
        process_end_time = time.time()
        print(f"[{datetime.now()}] /process endpoint finished in {process_end_time - process_start_time:.2f} seconds.")
        return ProcessResponse(
            session_id=session_id,
            message="Content processed successfully.",
            filenames=processed_files,
            processed_url=url
        )
    except Exception as e:
        print(f"[{datetime.now()}] ERROR in /process endpoint: {e}")
        raise HTTPException(status_code=500, detail=f"An internal error occurred: {e}")
    finally:
        print(f"[{datetime.now()}] Cleaning up temporary directory: {temp_dir}")
        shutil.rmtree(temp_dir)




# Combined api for inference 

class InferRequest(BaseModel):
    question: str = Field(..., description="User question / prompt")

    # Knowledge source (one of them)
    url: Optional[str] = None
    files: Optional[List[str]] = None  # future-safe

    # Optional controls
    voice_preset: Optional[str] = Field(None, description="Voice preset for TTS")
    return_audio: bool = Field(False, description="Return audio base64 or not")
    return_csv: bool = Field(False, description="Return blendshapes as CSV")


class InferResponse(BaseModel):
    answer: str
    blendshapes: List[dict]
    mapping: List[str] = Field(default_factory=list, description="Blendshape index-to-name mapping")
    frame_rate: int = Field(default=60, description="Blendshape frame rate")
    audio_base64: Optional[str] = None
    csv: Optional[str] = None

import json


import json


# =====================================================================================
# 8. Streaming WebSocket Endpoints
# =====================================================================================

# Streaming imports
from streaming.blendshape_worker import BlendshapeWorker
from streaming.kyutai_coordinator import KyutaiStreamCoordinator
from streaming.qwen_tts_worker import QwenTTSWorker
from streaming.optimized_blendshape_worker import OptimizedBlendshapeWorker
from streaming.performance_monitor import get_monitor


@app.websocket("/ws/infer/kyutai")
async def websocket_infer_kyutai(websocket: WebSocket):
    """
    Real-time streaming inference with PCM16 audio chunks.
    
    Features:
    - Real-time PCM16 audio streaming (configurable chunk_ms, default 50ms)
    - Joint audio-visual modeling with controlled delay
    - Adaptive buffering based on network conditions
    - Better synchronization between modalities
    - Graceful error recovery
    - Performance monitoring
    
    Client flow:
      1. Connect to ws://.../ws/infer/kyutai
      2. Send: {"type": "start", "session_id": "...", "question": "...", "chunk_ms": 50, "tts_instruct": "warm, gentle, soothing tone"}
      3. Receive progressive: text_chunk, audio_chunk (PCM16), blendshapes, status
      4. Send: {"type": "interrupt"} to stop
      5. Send: {"type": "buffer_adjust", "target_size": 3} to adjust buffering
    
    Audio format:
      - audio_chunk messages contain "audio_bytes_base64" (PCM16 little-endian)
      - "audio_format": "pcm_s16le", "channels": 1, "sample_rate": 24000
      - Client assembles PCM chunks into final WAV file
      - Use test_streaming_pcm.py as reference client
    """
    await websocket.accept()
    monitor = get_monitor()
    monitor.reset()
    
    try:
        init_msg = await websocket.receive_json()
        
        if init_msg.get("type") != "start":
            await websocket.send_json({
                "type": "status",
                "status": "error",
                "message": "First message must be type 'start'",
            })
            await websocket.close()
            return
        
        session_id = init_msg.get("session_id")
        question = init_msg.get("question")
        voice_preset = init_msg.get("voice_preset")
        tts_instruct = init_msg.get("tts_instruct")
        return_audio = init_msg.get("return_audio", True)
        use_qwen = init_msg.get("use_qwen", True)  # Default to Qwen3-TTS
        use_optimized_bs = init_msg.get("use_optimized_bs", True)
        chunk_ms = init_msg.get("chunk_ms")
        
        if not session_id or not question:
            await websocket.send_json({
                "type": "status",
                "status": "error",
                "message": "session_id and question are required",
            })
            await websocket.close()
            return
        
        chain = conversations.get(session_id)
        if not chain:
            await websocket.send_json({
                "type": "status",
                "status": "error",
                "message": "Session not found. Call /process first.",
            })
            await websocket.close()
            return
        
        # Create workers based on configuration
        device_str = "cuda" if torch.cuda.is_available() else "cpu"
        print(f"[{datetime.now()}] Using Qwen3-TTS worker")
        global _tts_reference_audio_path, _tts_reference_text
        tts_ref_audio = _tts_reference_audio_path
        tts_ref_text = _tts_reference_text

        if not tts_ref_audio or not tts_ref_text:
            await websocket.send_json(
                {
                    "type": "status",
                    "status": "error",
                    "message": "Missing voice-clone reference. Set TTS_REF_AUDIO_PATH and TTS_REF_TEXT on the server (RunPod) before using Base voice-clone streaming.",
                }
            )
            await websocket.close()
            return

        tts_worker = QwenTTSWorker(
            device=device_str,
            use_qwen3=True,
            reference_audio_path=tts_ref_audio,
            reference_text=tts_ref_text,
        )
        
        if use_optimized_bs:
            print(f"[{datetime.now()}] Using optimized blendshape worker")
            bs_worker = OptimizedBlendshapeWorker(blendshape_model, device, config)
        else:
            print(f"[{datetime.now()}] Using standard blendshape worker")
            bs_worker = BlendshapeWorker(blendshape_model, device, config)
        
        # Create Kyutai coordinator
        coordinator = KyutaiStreamCoordinator(
            websocket=websocket,
            tts_worker=tts_worker,
            blendshape_worker=bs_worker,
            config=config,
        )
        
        await coordinator.run_streaming_pipeline(
            rag_chain=chain,
            question=question,
            voice_preset=voice_preset,
            tts_instruct=tts_instruct,
            return_audio=return_audio,
            chunk_ms=chunk_ms,
        )
        
        # Print performance summary
        monitor.print_summary()
    
    except WebSocketDisconnect:
        print(f"[{datetime.now()}] WebSocket client disconnected")
    except Exception as e:
        print(f"[{datetime.now()}] WebSocket error: {e}")
        try:
            await websocket.send_json({
                "type": "status",
                "status": "error",
                "message": str(e),
            })
        except Exception:
            pass
    finally:
        try:
            await websocket.close()
        except Exception:
            pass


@app.get("/performance/summary")
async def get_performance_summary():
    """Get current performance metrics."""
    monitor = get_monitor()
    return monitor.get_summary()

# flask_Character.py
# Drop-in replacement. DO NOT keep any other conflicting copies with a different name.

from fastapi import FastAPI, Request, File, UploadFile, Form, HTTPException, BackgroundTasks
from fastapi.responses import JSONResponse, PlainTextResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field
from typing import List, Optional, Dict, Any
from datetime import datetime
import time
import io
import base64
import os
import uuid
import shutil
import zipfile
import requests
import tempfile

# ML / utils imports
import torch
import numpy as np
import pandas as pd
import docx
from pptx import Presentation
from PyPDF2 import PdfReader
import scipy.io.wavfile as wavfile
from transformers import AutoProcessor, BarkModel

# Langchain & vectorstore
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from langchain_groq import ChatGroq
from langchain.prompts import PromptTemplate

# Local helper functions (these files should exist)
from utils.generate_face_shapes import generate_facial_data_from_bytes
from utils.config import config
from utils.model.model import load_model  # we will call this ONLY inside lazy loader

# Safety envs
os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")
os.environ.setdefault("TRANSFORMERS_CACHE", "/app/.cache/huggingface")
os.environ.setdefault("HF_HOME", "/app/.cache/huggingface")
os.environ.setdefault("TORCH_HOME", "/app/.cache/torch")

# App init
app = FastAPI(
    title="Intelligent Document & Web API",
    description="RAG + Bark TTS + blendshape from audio (serverless-safe).",
    version="2.0.1",
)

# Static audio directory
STATIC_AUDIO_DIR = "generated_audio"
os.makedirs(STATIC_AUDIO_DIR, exist_ok=True)
app.mount("/audio", StaticFiles(directory=STATIC_AUDIO_DIR), name="audio")

# -------------------------
# Global state (lazy / startup loaded)
# -------------------------
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "")
if not GROQ_API_KEY:
    # keep service startable but warn
    print("[WARN] GROQ_API_KEY not set; ChatGroq calls will fail until set.")

conversations: Dict[str, ConversationalRetrievalChain] = {}
embeddings_model = None

# Bark TTS global placeholders (loaded at startup)
bark_processor = None
bark_model = None
bark_sr = 24000
bark_device = None
bark_available = False

# Blendshape (lazy)
_blendshape_model = None
_blendshape_device = None

# -------------------------
# Prompt template (unchanged)
# -------------------------
PROFESSIONAL_RAG_PROMPT_TEMPLATE = """
You are a highly intelligent and diligent AI research assistant. Your primary goal is to provide accurate, concise, and helpful answers based *only* on the context provided.

**Instructions:**
1.  **Analyze the Context:** Carefully read and understand the `context` provided below. It is your only source of truth.
2.  **Answer the Question:** Use the context to answer the user's `question`.
3.  **Strict Grounding:** Do not use any external knowledge. If the answer is not in the context, you MUST state: "I am sorry, but the information required to answer your question is not available in the provided documents." Do not try to guess or infer information that isn't explicitly stated.
4.  **Synthesize Information:** If the question requires combining information from multiple parts of the context, synthesize a coherent answer.
5.  **Clarity and Conciseness:** Provide a clear and direct answer. If appropriate, use bullet points to structure complex information.
6.  **Cite Sources (if applicable):** While not strictly required, if you can identify the source of a piece of information within the context, it's good practice to mention it.

**Context:**
{context}

**Chat History:**
{chat_history}

**Question:**
{question}

**Answer:**
"""
RAG_PROMPT = PromptTemplate.from_template(PROFESSIONAL_RAG_PROMPT_TEMPLATE)

# -------------------------
# Lazy loader for blendshape model (serverless safe)
# -------------------------
def get_blendshape_model():
    """
    Lazy-load the blendshape model. This function will be called on the first request
    to /audio_to_blendshapes. It wraps torch.load temporarily to ensure compatibility
    with weights_only behavior in PyTorch 2.6+ and restores the original torch.load.
    """
    global _blendshape_model, _blendshape_device

    if _blendshape_model is not None:
        return _blendshape_model, _blendshape_device

    import torch as _torch
    from utils.model.model import load_model as _loader

    device = _torch.device("cuda" if _torch.cuda.is_available() else "cpu")
    model_path = "utils/model/model.pth"

    print(f"[{datetime.now()}] Loading blendshape model on device: {device} (path={model_path})")

    # Save original torch.load and wrap it to ensure weights_only=False unless explicitly set
    orig_torch_load = _torch.load

    def torch_load_wrapper(f, *args, **kwargs):
        if "weights_only" not in kwargs:
            kwargs["weights_only"] = False
        return orig_torch_load(f, *args, **kwargs)

    _torch.load = torch_load_wrapper
    try:
        _blendshape_model = _loader(model_path, config, device)
        _blendshape_device = device
        print(f"[{datetime.now()}] Blendshape model loaded successfully.")
    except Exception as ex:
        # restore original before raising and log error
        _torch.load = orig_torch_load
        print(f"[{datetime.now()}] ERROR loading blendshape model: {repr(ex)}")
        raise
    finally:
        # Ensure restore in all cases
        _torch.load = orig_torch_load

    return _blendshape_model, _blendshape_device

# -------------------------
# Startup: load embeddings + Bark TTS (safe to run at startup)
# -------------------------
@app.on_event("startup")
async def load_startup_models():
    """
    Load embeddings and Bark TTS at startup (so errors are visible early).
    This avoids loading the blendshape model at import time.
    """
    global embeddings_model, bark_processor, bark_model, bark_sr, bark_device, bark_available

    print(f"[{datetime.now()}] Startup: loading embeddings and Bark (if available)...")

    # Load embeddings (best-effort)
    try:
        embeddings_model = HuggingFaceEmbeddings(
            model_name="all-MiniLM-L6-v2",
            model_kwargs={"device": "cuda" if torch.cuda.is_available() else "cpu"},
        )
        print(f"[{datetime.now()}] Embeddings model loaded.")
    except Exception as e:
        embeddings_model = None
        print(f"[{datetime.now()}] Failed to load embeddings model: {e}")

    # Load Bark (best-effort)
    try:
        device_str = "cuda" if torch.cuda.is_available() else "cpu"
        device = torch.device(device_str)

        bark_processor = AutoProcessor.from_pretrained("suno/bark")
        bark_model = BarkModel.from_pretrained("suno/bark").to(device)

        bark_sr = getattr(bark_model, "generation_config", {}).sample_rate if hasattr(bark_model, "generation_config") else 24000
        bark_device = device
        bark_available = True
        print(f"[{datetime.now()}] Bark loaded on {device} (sr={bark_sr}).")
    except Exception as e:
        bark_available = False
        print(f"[{datetime.now()}] Bark load failed; TTS disabled: {e}")

    print(f"[{datetime.now()}] Startup complete.")

# -------------------------
# Helper: Content extraction
# -------------------------
class ContentExtractor:
    def from_url(self, url: str) -> str:
        try:
            response = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
            response.raise_for_status()
            soup = BeautifulSoup(response.content, "html.parser")
            for element in soup(["script", "style", "header", "footer", "nav", "aside"]):
                element.decompose()
            return soup.get_text(separator="\n", strip=True)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Failed to fetch URL: {e}")

    def from_pdf(self, file_stream) -> str:
        try:
            reader = PdfReader(file_stream)
            return "".join(page.extract_text() or "" for page in reader.pages)
        except Exception as e:
            return f"[PDF Parsing Error: {e}]"

    def from_docx(self, file_stream) -> str:
        try:
            doc = docx.Document(file_stream)
            return "\n".join(para.text for para in doc.paragraphs if para.text)
        except Exception as e:
            return f"[DOCX Parsing Error: {e}]"

    def from_pptx(self, file_stream) -> str:
        try:
            prs = Presentation(file_stream)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception as e:
            return f"[PPTX Parsing Error: {e}]"

    def from_excel(self, file_stream) -> str:
        try:
            xls = pd.ExcelFile(file_stream)
            text = ""
            for sheet in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet)
                text += f"--- Sheet: {sheet} ---\n{df.to_string()}\n\n"
            return text
        except Exception as e:
            return f"[Excel Parsing Error: {e}]"

    def from_zip(self, file_path, temp_dir) -> str:
        text = ""
        try:
            with zipfile.ZipFile(file_path) as z:
                z.extractall(path=temp_dir)
                for root, _, files in os.walk(temp_dir):
                    for fname in files:
                        fpath = os.path.join(root, fname)
                        text += self.from_file_path(fpath) + "\n\n"
        except Exception as e:
            text += f"[ZIP Parsing Error: {e}]"
        return text

    def from_file_path(self, file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        with open(file_path, "rb") as f:
            if ext == ".pdf":
                return self.from_pdf(f)
            elif ext == ".docx":
                return self.from_docx(f)
            elif ext == ".pptx":
                return self.from_pptx(f)
            elif ext in (".xls", ".xlsx"):
                return self.from_excel(f)
            else:
                return f"[Unsupported file type: {ext}]"

# -------------------------
# RAG helper
# -------------------------
def get_rag_chain(text_chunks: List[str]) -> ConversationalRetrievalChain:
    if not text_chunks:
        raise ValueError("No text chunks supplied.")
    if embeddings_model is None:
        raise RuntimeError("Embeddings model not available.")
    vector_store = FAISS.from_texts(texts=text_chunks, embedding=embeddings_model)
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True, output_key="answer")
    llm = ChatGroq(
        model="meta-llama/llama-4-scout-17b-16e-instruct",
        temperature=0.7,
        max_tokens=800,
        groq_api_key=GROQ_API_KEY
    )
    rag_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vector_store.as_retriever(),
        memory=memory,
        return_source_documents=True,
        combine_docs_chain_kwargs={"prompt": RAG_PROMPT},
    )
    return rag_chain

from collections.abc import Mapping

# -------------------------
# TTS: in-memory generation helper (uses startup-loaded bark variables)
# -------------------------
def generate_speech_content_base64(text: str, voice_preset: Optional[str]) -> Optional[str]:
    if not bark_available or not text:
        return None

    # prepare inputs
    attempts = [
        {"voice_preset": voice_preset, "return_tensors": "pt", "padding": True},
        {"voice_preset": voice_preset, "return_tensors": "pt"},
        {"return_tensors": "pt", "padding": True},
        {"return_tensors": "pt"},
        {}
    ]
    inputs = None
    last_exc = None
    for kw in attempts:
        try:
            safe_kw = {k: v for k, v in kw.items() if v is not None}
            inputs = bark_processor(text, **safe_kw)
            break
        except Exception as e:
            last_exc = e

    if inputs is None:
        # fallback to tokenizer if available
        if hasattr(bark_processor, "tokenizer"):
            try:
                inputs = bark_processor.tokenizer(text, return_tensors="pt", padding=True)
            except Exception as e:
                return None
        else:
            return None

    # safe move to device
    try:
        if hasattr(inputs, "to") and callable(inputs.to):
            inputs = inputs.to(bark_device)
    except Exception:
        pass

    # recursive conversion (simple version)
    def to_tensors(obj):
        if torch.is_tensor(obj):
            return obj.to(bark_device)
        if isinstance(obj, np.ndarray):
            return torch.tensor(obj, device=bark_device)
        if isinstance(obj, Mapping):
            return {k: to_tensors(v) for k, v in obj.items()}
        if isinstance(obj, (list, tuple)):
            return type(obj)(to_tensors(v) for v in obj)
        return obj

    converted = to_tensors(inputs)

    # ensure attention mask
    if isinstance(converted, dict) and "input_ids" in converted and "attention_mask" not in converted:
        try:
            input_ids = converted["input_ids"]
            pad_id = getattr(bark_model.config, "pad_token_id", None)
            if pad_id is not None:
                mask = (input_ids != int(pad_id)).long().to(bark_device)
            else:
                mask = torch.ones_like(input_ids, dtype=torch.long, device=bark_device)
            converted["attention_mask"] = mask
        except Exception:
            pass

    # generate
    try:
        bark_model.eval()
        with torch.no_grad():
            gen_kwargs = {}
            try:
                if hasattr(bark_model, "generate"):
                    output = bark_model.generate(**converted, **gen_kwargs)
                else:
                    return None
            except TypeError:
                output = bark_model.generate(**converted)

        # extract waveform
        audio_array = None
        if isinstance(output, dict):
            for k in ("audio", "audios", "waveform", "wav", "output_audio"):
                if k in output:
                    audio_array = output[k]
                    break
            if audio_array is None and "outputs" in output:
                cand = output["outputs"]
                if isinstance(cand, (list, tuple)) and cand:
                    audio_array = cand[0]
        elif isinstance(output, (list, tuple)):
            audio_array = output[0]
        else:
            audio_array = output

        if audio_array is None:
            return None

        if hasattr(audio_array, "cpu"):
            audio_np = audio_array.cpu().numpy().squeeze()
        else:
            audio_np = np.asarray(audio_array).squeeze()

        if audio_np.ndim > 1:
            audio_np = audio_np.mean(axis=-1)

        sr = int(bark_sr) if bark_sr else 24000
        if np.issubdtype(audio_np.dtype, np.floating):
            audio_np = np.clip(audio_np, -1.0, 1.0)
            audio_int16 = (audio_np * 32767.0).astype(np.int16)
        else:
            audio_int16 = audio_np.astype(np.int16)

        buf = io.BytesIO()
        wavfile.write(buf, sr, audio_int16)
        buf.seek(0)
        bts = buf.read()
        return base64.b64encode(bts).decode("utf-8")

    except Exception:
        return None

# -------------------------
# Pydantic models
# -------------------------
class ProcessResponse(BaseModel):
    session_id: str
    message: str
    filenames: List[str]
    processed_url: Optional[str] = None

class QueryRequest(BaseModel):
    session_id: str
    question: str
    voice_preset: Optional[str] = None

class SourceDocument(BaseModel):
    page_content: str
    metadata: Dict[str, Any]

class QueryResponse(BaseModel):
    answer: str
    source_documents: List[SourceDocument]
    audio_base64: Optional[str] = None

# -------------------------
# API endpoints
# -------------------------
@app.get("/")
def health():
    return {"status": "ok", "time": datetime.utcnow().isoformat()}

@app.post("/audio_to_blendshapes")
async def audio_to_blendshapes_route(request: Request):
    audio_bytes = await request.body()
    if not audio_bytes:
        raise HTTPException(status_code=400, detail="No audio bytes received.")

    # lazy load blendshape model here
    blendshape_model, device = get_blendshape_model()

    # call your existing generator (ensure generate_facial_data_from_bytes matches signature)
    generated_facial_data = generate_facial_data_from_bytes(audio_bytes, blendshape_model, device, config)
    if isinstance(generated_facial_data, np.ndarray):
        generated_facial_data = generated_facial_data.tolist()

    return JSONResponse(content={"blendshapes": generated_facial_data})

@app.post("/process", response_model=ProcessResponse)
async def process_content(
    files: Optional[List[UploadFile]] = File(None),
    url: Optional[str] = Form(None),
    prompt_template: Optional[str] = Form(None)
):
    if not files and not url:
        raise HTTPException(status_code=400, detail="Provide files or a url.")

    extractor = ContentExtractor()
    tmpdir = f"temp_{uuid.uuid4().hex}"
    os.makedirs(tmpdir, exist_ok=True)
    raw_text = ""
    processed_files = []

    try:
        if url:
            raw_text += extractor.from_url(url) + "\n\n"

        if files:
            for f in files:
                safe_name = os.path.basename(f.filename)
                path = os.path.join(tmpdir, safe_name)
                with open(path, "wb") as out:
                    shutil.copyfileobj(f.file, out)
                ext = os.path.splitext(safe_name)[1].lower()
                if ext == ".zip":
                    raw_text += extractor.from_zip(path, os.path.join(tmpdir, "unzipped"))
                else:
                    raw_text += extractor.from_file_path(path) + "\n\n"
                processed_files.append(safe_name)

        if not raw_text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted.")

        splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200, length_function=len)
        text_chunks = splitter.split_text(raw_text)
        session_id = str(uuid.uuid4())
        conversations[session_id] = get_rag_chain(text_chunks)

        if prompt_template:
            conversations[session_id].combine_docs_chain.llm_chain.prompt = PromptTemplate.from_template(prompt_template)

        return ProcessResponse(session_id=session_id, message="Processed", filenames=processed_files, processed_url=url)
    finally:
        try:
            shutil.rmtree(tmpdir)
        except Exception:
            pass

@app.post("/query", response_model=QueryResponse)
async def query_session(request: QueryRequest, background_tasks: BackgroundTasks):
    chain = conversations.get(request.session_id)
    if not chain:
        raise HTTPException(status_code=404, detail="Session not found")

    result = chain.invoke({"question": request.question})
    answer = result.get("answer", "No answer generated.")
    source_docs = result.get("source_documents", [])
    validated = [SourceDocument(page_content=getattr(d, "page_content", ""), metadata=getattr(d, "metadata", {})) for d in source_docs]

    audio_b64 = None
    if bark_available:
        audio_b64 = generate_speech_content_base64(answer, voice_preset=request.voice_preset)

    return QueryResponse(answer=answer, source_documents=validated, audio_base64=audio_b64)
